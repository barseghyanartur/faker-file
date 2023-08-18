from io import BytesIO
from typing import Callable, Optional, Union, overload

from faker import Faker
from faker.generator import Generator
from faker.providers import BaseProvider
from faker.providers.python import Provider
from pptx import Presentation
from pptx.util import Inches

from ..base import DEFAULT_FORMAT_FUNC, BytesValue, FileMixin, StringValue
from ..constants import DEFAULT_TEXT_MAX_NB_CHARS
from ..registry import FILE_REGISTRY
from ..storages.base import BaseStorage
from ..storages.filesystem import FileSystemStorage

__author__ = "Artur Barseghyan <artur.barseghyan@gmail.com>"
__copyright__ = "2022-2023 Artur Barseghyan"
__license__ = "MIT"
__all__ = ("PptxFileProvider",)


class PptxFileProvider(BaseProvider, FileMixin):
    """PPTX file provider.

    Usage example:

    .. code-block:: python

        from faker import Faker
        from faker_file.providers.pptx_file import PptxFileProvider

        FAKER = Faker()
        FAKER.add_provider(PptxFileProvider)

        file = FAKER.pptx_file()

    Usage example with options:

    .. code-block:: python

        file = FAKER.pptx_file(
            prefix="zzz",
            max_nb_chars=100_000,
            wrap_chars_after=80,
        )

    Usage example with `FileSystemStorage` storage (for `Django`):

    .. code-block:: python

        from django.conf import settings
        from faker_file.storages.filesystem import FileSystemStorage

        file = FAKER.pptx_file(
            storage=FileSystemStorage(
                root_path=settings.MEDIA_ROOT,
                rel_path="tmp",
            ),
            prefix="zzz",
            max_nb_chars=100_000,
            wrap_chars_after=80,
        )
    """

    extension: str = "pptx"

    @overload
    def pptx_file(
        self: "PptxFileProvider",
        storage: Optional[BaseStorage] = None,
        basename: Optional[str] = None,
        prefix: Optional[str] = None,
        max_nb_chars: int = DEFAULT_TEXT_MAX_NB_CHARS,
        wrap_chars_after: Optional[int] = None,
        content: Optional[str] = None,
        format_func: Callable[
            [Union[Faker, Generator, Provider], str], str
        ] = DEFAULT_FORMAT_FUNC,
        raw: bool = True,
        **kwargs,
    ) -> BytesValue:
        ...

    @overload
    def pptx_file(
        self: "PptxFileProvider",
        storage: Optional[BaseStorage] = None,
        basename: Optional[str] = None,
        prefix: Optional[str] = None,
        max_nb_chars: int = DEFAULT_TEXT_MAX_NB_CHARS,
        wrap_chars_after: Optional[int] = None,
        content: Optional[str] = None,
        format_func: Callable[
            [Union[Faker, Generator, Provider], str], str
        ] = DEFAULT_FORMAT_FUNC,
        **kwargs,
    ) -> StringValue:
        ...

    def pptx_file(
        self: "PptxFileProvider",
        storage: Optional[BaseStorage] = None,
        basename: Optional[str] = None,
        prefix: Optional[str] = None,
        max_nb_chars: int = DEFAULT_TEXT_MAX_NB_CHARS,
        wrap_chars_after: Optional[int] = None,
        content: Optional[str] = None,
        format_func: Callable[
            [Union[Faker, Generator, Provider], str], str
        ] = DEFAULT_FORMAT_FUNC,
        raw: bool = False,
        **kwargs,
    ) -> Union[BytesValue, StringValue]:
        """Generate a file with random text.

        :param storage: Storage. Defaults to `FileSystemStorage`.
        :param basename: File basename (without extension).
        :param prefix: File name prefix.
        :param max_nb_chars: Max number of chars for the content.
        :param wrap_chars_after: If given, the output string would be separated
             by line breaks after the given position.
        :param content: File content. Might contain dynamic elements, which
            are then replaced by correspondent fixtures.
        :param format_func: Callable responsible for formatting template
            strings.
        :param raw: If set to True, return `BytesValue` (binary content of
            the file). Otherwise, return `StringValue` (path to the saved
            file).
        :return: Relative path (from root directory) of the generated file
            or raw content of the file.
        """
        # Generic
        if storage is None:
            storage = FileSystemStorage()

        filename = storage.generate_filename(
            extension=self.extension,
            prefix=prefix,
            basename=basename,
        )

        content = self._generate_text_content(
            max_nb_chars=max_nb_chars,
            wrap_chars_after=wrap_chars_after,
            content=content,
            format_func=format_func,
        )

        data = {"content": content, "filename": filename, "storage": storage}

        stream = BytesIO()
        prs = Presentation()
        prs.slide_width = Inches(25)
        prs.slide_height = Inches(16)
        # Make a blank slide with a text box with random text
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_box = slide.shapes.add_textbox(0, 0, Inches(25), Inches(16))
        text_box.text = content
        text_box.text_frame.word_wrap = True
        prs.save(stream)

        if raw:
            raw_content = BytesValue(stream.getvalue())
            raw_content.data = data
            return raw_content

        storage.write_bytes(filename, stream.getvalue())

        # Generic
        file_name = StringValue(storage.relpath(filename))
        file_name.data = data
        FILE_REGISTRY.add(file_name)
        return file_name
